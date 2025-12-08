import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Optional

# import the python-pptx and PyMuPDF library
from pptx import Presentation
import fitz

# JSON builder
def build_json(
    *,
    course_id: str,
    lecture_title: str,
    source_file: str,
    source_type: str,
    slides: List[Dict[str, Any]],
) -> Dict[str, Any]:
    return {
        "course_id": course_id,         # ex: CSE368
        "lecture_title": lecture_title,
        "source_file": source_file,     # the file
        "source_type": source_type,     # "pptx" or "pdf"
        "slides": slides                # list of slide/page dicts
    }

def build_slide_obj(
    *,
    index: int,
    title: str,
    text_blocks: List[str],
    notes: str = "",
    images: Optional[List[str]] = None,
) -> Dict[str, Any]:
    if images is None:
        images = []
    return {
        "index": index,             # slide/page number starting at 1
        "title": title,             # title or heading
        "text_blocks": text_blocks, # text chunks
        "notes": notes,             # speaker notes
        "images": images,           # images (for future if we have time)
    }

def extract_pptx(pptx_path: Path) -> List[Dict[str, Any]]:

    ppt = Presentation(pptx_path)
    slides: List[Dict[str, Any]] = []

    for i, slide in enumerate(ppt.slides, start = 1):
        title = ""
        text_blocks: List[str] = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            raw = shape.text or ""
            text = raw.strip()
            if not text:
                continue
            
            if title == "":
                title = text
            else:
                text_blocks.append(text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes_raw = slide.notes_slide.notes_text_frame.text or ""
            notes = notes_raw.strip()

        slide_obj = build_slide_obj(
            index = i, 
            title = title if title else f"slide {i}",
            text_blocks = text_blocks,
            notes = notes,
            images = [],            # if we have time to impliment
        )
        slides.append(slide_obj)
            
    return slides

def extract_pdf(pdf_path: Path) -> List[Dict[str, Any]]:
    
    pdfdoc = fitz.open(pdf_path)
    slides: List[Dict[str, Any]] = []

    for i, page in enumerate(pdfdoc, start = 1):
        blocks = page.get_text("blocks")
        
        page_blocks: List[str] = []

        for block in blocks:
            text = (block[4] or "").strip()
            if text:
                page_blocks.append(text)

        if not page_blocks:
            slide_obj = build_slide_obj(
                index = i,
                title = f"Page {i}",
                text_blocks = [],
            )
            slides.append(slide_obj)
            continue

        first_block = page_blocks[0]
        if len(first_block.splitlines()) <= 3:
            title = first_block
            text_blocks = page_blocks[1:]
        else:
            title = f"Page {i}"
            text_blocks = page_blocks

        slide_obj = build_slide_obj(
            index = i,
            title = title,
            text_blocks = text_blocks,
        )
        slides.append(slide_obj)
    
    return slides

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description = "Extract PPTX/PDF content into JSON.")
    parser.add_argument("--input", required = True, type = Path, help = "Path to input pptx/pdf")
    parser.add_argument("--course-id", required = True, help = "Course identifier")
    parser.add_argument("--lecture-title", required = True, help = "Lecture title")
    parser.add_argument("--output", type = Path, help = "Write JSON to this path")
    return parser.parse_args()

def main() -> None:
    args = parse_args()
    input_path: Path = args.input

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    ext = input_path.suffix.lower()
    if ext == ".pptx":
        slides = extract_pptx(input_path)
        source_type = "pptx"
    elif ext == ".pdf":
        slides = extract_pdf(input_path)
        source_type = "pdf"
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    doc_json = build_json(
        course_id = args.course_id,
        lecture_title = args.lecture_title,
        source_file = str(input_path),
        source_type = source_type,
        slides = slides,
    )

    output_json = json.dumps(doc_json, indent = 2)
    if args.output:
        args.output.write_text(output_json, encoding = "utf-8")
    else:
        print(output_json)

if __name__ == "__main__":
    main()